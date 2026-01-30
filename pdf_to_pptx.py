from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os
import sys

def convert_pdf_to_pptx(pdf_path, output_path=None):
    """
    Converts a PDF to a PPTX by rendering each page as an image 
    and placing it onto a slide.
    """
    if not output_path:
        output_path = os.path.splitext(pdf_path)[0] + ".pptx"

    print(f"--- Converting to PPTX: {os.path.basename(pdf_path)} ---")

    try:
        # 1. Convert PDF pages to Images
        print("  • Rendering PDF pages to images (this may take a moment)...")
        images = convert_from_path(pdf_path, dpi=200)

        # 2. Create PPTX
        prs = Presentation()
        
        # Set slide size to match the first page image aspect ratio if possible
        # Default PPTX is 16:9 or 4:3. Here we just fit image to slide.
        
        # Standard 16:9 dimensions in EMU (English Metric Units)
        # 1 inch = 914400 EMU
        # Default slide width 10 inches, height 7.5 inches (4:3) or 13.333x7.5 (16:9)
        # We will adjust slide size to match image size for better quality
        
        if images:
            width, height = images[0].size
            # pptx uses EMU. Let's approximate.
            # 72 dpi is standard screen? pdf2image uses 200 dpi.
            # To avoid huge slides, we scale down or just set slide size to match image aspect.
            # Actually, easiest is: set slide layout to blank, add picture covering whole slide.
            
            prs.slide_width = int(width * 914400 / 200) # roughly convert px to emu based on dpi
            prs.slide_height = int(height * 914400 / 200)

        for i, img in enumerate(images):
            # Save temp image
            temp_img_path = f"temp_slide_{i}.jpg"
            img.save(temp_img_path, 'JPEG')

            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)

            # Add image to slide
            slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # Clean up temp image
            os.remove(temp_img_path)
            print(f"  ✓ Processed Page {i+1}")

        prs.save(output_path)
        print(f"✓ Saved to: {output_path}")
        return output_path

    except Exception as e:
        print(f"✗ Failed to convert {pdf_path}: {e}")
        return None

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pdf_to_pptx.py <file.pdf>")
    else:
        convert_pdf_to_pptx(sys.argv[1])
